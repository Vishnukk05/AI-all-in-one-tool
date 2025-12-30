import os
import logging
import datetime
import uuid
import time
import psutil
import threading
from io import BytesIO

# --- FLASK IMPORTS ---
from flask import Flask, render_template, request, jsonify, Response, session
from dotenv import load_dotenv

# --- AI & MEDIA IMPORTS ---
from groq import Groq
from gtts import gTTS
from xhtml2pdf import pisa
from pptx import Presentation
import speech_recognition as sr
import PIL.Image
from moviepy.video.io.VideoFileClip import VideoFileClip

# --- SETUP ENV ---
basedir = os.path.abspath(os.path.dirname(__file__))
load_dotenv(os.path.join(basedir, '.env'))

app = Flask(__name__)
# CRITICAL: Change this key for production security
app.secret_key = os.environ.get("SECRET_KEY", "super_secret_key_123")

# --- CONFIGURATION ---
API_KEY = os.environ.get("GROQ_API_KEY")
STATIC_FOLDER = os.path.join(basedir, 'static')
if not os.path.exists(STATIC_FOLDER):
    os.makedirs(STATIC_FOLDER)

# --- ADMIN CONFIG ---
ADMIN_USER = os.environ.get("ADMIN_USER", "Admin")
ADMIN_PASS = os.environ.get("ADMIN_PASS", "admin123") 

# --- STATS TRACKING ---
global_stats = {
    "text_gen": 0, "audio_gen": 0, "transcribe": 0, "pdf_gen": 0, 
    "chat_msgs": 0, "code_review": 0, "quiz_gen": 0,
    "file_conv": 0, "compression": 0, "vid_audio": 0
}

def increment_stat(field):
    if field in global_stats:
        global_stats[field] += 1

# --- UTILITY: FILE CLEANUP ---
def cleanup_old_files():
    """Deletes files in static folder older than 30 minutes to save space."""
    try:
        now = time.time()
        for f in os.listdir(STATIC_FOLDER):
            fpath = os.path.join(STATIC_FOLDER, f)
            if os.path.isfile(fpath):
                # If file is older than 30 mins (1800 seconds)
                if now - os.path.getmtime(fpath) > 1800:
                    try: os.remove(fpath)
                    except: pass
    except Exception as e:
        print(f"Cleanup Error: {e}")

# Run cleanup in background on every request
@app.before_request
def before_request_cleanup():
    if request.endpoint != 'static': 
        threading.Thread(target=cleanup_old_files).start()

# --- UTILITY: CLEAN AI TEXT ---
def clean_ai_text(text):
    """Removes Markdown fences like ```html from AI response."""
    if not text: return ""
    return text.replace("```html", "").replace("```json", "").replace("```", "").strip()

# --- UTILITY: AI WRAPPER ---
def get_groq_response(system_prompt, user_prompt, temperature=0.5):
    if not API_KEY:
        print("❌ Error: API Key is missing.")
        return None
    try:
        client = Groq(api_key=API_KEY)
        completion = client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt}
            ],
            temperature=temperature,
            max_tokens=2048
        )
        return completion.choices[0].message.content
    except Exception as e:
        print(f"❌ Groq API Error: {e}")
        return None

# ==============================================================================
#                               CORE ROUTES
# ==============================================================================

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/health')
def health():
    return "OK", 200

# --- AUTH ---
@app.route('/login', methods=['POST'])
def login():
    data = request.json
    if data.get('username') == ADMIN_USER and data.get('password') == ADMIN_PASS:
        session['is_admin'] = True
        return jsonify({"success": True})
    return jsonify({"success": False, "error": "Invalid Credentials"}), 401

@app.route('/logout', methods=['POST'])
def logout():
    session.pop('is_admin', None)
    return jsonify({"success": True})

@app.route('/check-auth')
def check_auth():
    return jsonify({"is_admin": session.get('is_admin', False)})

# --- STATS & REPORTS (FIXED: CPU Load) ---
@app.route('/api/stats')
def get_stats():
    cpu, ram = 0, 0
    # Only calculate stats if admin is logged in (saves resources)
    if session.get('is_admin', False):
        try:
            # FIX: interval=0.1 forces psutil to measure over 0.1s
            # This fixes the "0.0%" issue
            cpu = psutil.cpu_percent(interval=0.1) 
            ram = psutil.virtual_memory().percent
        except: pass
    return jsonify({"cpu": cpu, "ram": ram, "usage": global_stats})

@app.route('/download-report')
def download_report():
    if not session.get('is_admin'):
        return "Unauthorized", 401
    
    now = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    report_lines = [
        "========================================",
        "       AI WORKSPACE SYSTEM REPORT       ",
        "========================================",
        f"Generated: {now}",
        f"Server CPU: {psutil.cpu_percent(interval=0.1)}%",
        f"Server RAM: {psutil.virtual_memory().percent}%",
        "",
        "----------- FEATURE USAGE ------------",
    ]
    
    total_ops = sum(global_stats.values())
    for k, v in global_stats.items():
        report_lines.append(f"{k.ljust(15)} : {v}")
    
    report_lines.append("--------------------------------------")
    report_lines.append(f"TOTAL OPERATIONS : {total_ops}")
    report_lines.append("========================================")
    
    return Response(
        "\n".join(report_lines), 
        mimetype="text/plain", 
        headers={"Content-disposition": f"attachment; filename=System_Report.txt"}
    )

# ==============================================================================
#                               AI TOOLS
# ==============================================================================

@app.route('/chat', methods=['POST'])
def chat():
    increment_stat('chat_msgs')
    msg = request.form.get('message', '')
    if not msg: return jsonify({"success": False, "error": "Empty message"}), 400

    history = session.get('chat_history', [])
    sys_msg = {"role": "system", "content": "You are a helpful AI office assistant. Be concise."}
    messages = [sys_msg] + history + [{"role": "user", "content": msg}]
    
    try:
        client = Groq(api_key=API_KEY)
        completion = client.chat.completions.create(
            model="llama-3.3-70b-versatile", messages=messages, temperature=0.7
        )
        ai_reply = completion.choices[0].message.content
        
        history.append({"role": "user", "content": msg})
        history.append({"role": "assistant", "content": ai_reply})
        if len(history) > 6: history = history[-6:]
        session['chat_history'] = history
        
        return jsonify({"success": True, "response": ai_reply})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500

@app.route('/clear-chat', methods=['POST'])
def clear_chat():
    session.pop('chat_history', None)
    return jsonify({"success": True})

@app.route('/generate-minutes', methods=['POST'])
def generate_minutes():
    increment_stat('text_gen')
    notes = request.form.get('notes', '')
    res = get_groq_response(
        "You are a secretary. Convert raw notes into professional Meeting Minutes. Use Markdown.",
        f"Here are the notes:\n{notes}"
    )
    return jsonify({"success": True, "minutes": res if res else "AI Service Busy"})

@app.route('/generate-email', methods=['POST'])
def generate_email():
    increment_stat('text_gen')
    to = request.form.get('recipient', 'Team')
    topic = request.form.get('topic', 'Update')
    res = get_groq_response(
        "You are a professional email drafter.",
        f"Write an email to {to} about: {topic}"
    )
    return jsonify({"success": True, "email_content": res if res else "AI Service Busy"})

@app.route('/review-code', methods=['POST'])
def review_code():
    increment_stat('code_review')
    code = request.form.get('code', '')
    res = get_groq_response(
        "You are a Senior Developer. Review this code, find bugs, and suggest fixes. Use Markdown.",
        code
    )
    return jsonify({"success": True, "review": res if res else "AI Service Busy"})

@app.route('/translate', methods=['POST'])
def translate():
    increment_stat('text_gen')
    text = request.form.get('text', '')
    target = request.form.get('target_language', 'English')
    res = get_groq_response(
        "You are a professional translator. Output ONLY the translated text.",
        f"Translate this text to {target}:\n{text}"
    )
    return jsonify({"success": True, "translation": res if res else "Error"})

# --- QUIZ GENERATOR (FIXED: Alignment & PDF Style) ---
@app.route('/generate-quiz', methods=['POST'])
def generate_quiz():
    increment_stat('quiz_gen')
    topic = request.form.get('topic', 'General Knowledge')
    count = request.form.get('count', '5')
    
    # 1. Strict HTML Prompt to AI
    prompt = (
        f"Create a {count}-question Multiple Choice Quiz about '{topic}'.\n"
        "Output ONLY raw HTML content (no ```html fences).\n"
        "Use this EXACT structure for every question:\n"
        "<div class='question-box'>\n"
        "  <h3 class='q-title'>1. Question text here?</h3>\n"
        "  <ul class='options-list'>\n"
        "    <li>A) Option 1</li>\n"
        "    <li>B) Option 2</li>\n"
        "    <li>C) Option 3</li>\n"
        "    <li>D) Option 4</li>\n"
        "  </ul>\n"
        "</div>\n"
        "At the very end, add: <h4>Answer Key</h4>\n"
        "<table class='answer-key'>...</table>"
    )
    
    raw_res = get_groq_response("You are a strict HTML quiz generator.", prompt)
    if not raw_res: return jsonify({"success": False, "error": "AI Failed"})
    
    clean_html = clean_ai_text(raw_res)
    
    # 2. PDF CSS Wrapper (FIXED ALIGNMENT)
    pdf_html = f"""
    <html>
    <head>
        <style>
            @page {{ size: A4; margin: 2cm; }}
            body {{ font-family: Helvetica, sans-serif; font-size: 12px; color: #000; line-height: 1.4; }}
            
            /* Header */
            h1 {{ text-align: center; color: #4f46e5; border-bottom: 2px solid #eee; padding-bottom: 10px; margin-bottom: 20px; }}
            
            /* Questions */
            .question-box {{ margin-bottom: 10px; page-break-inside: avoid; }}
            .q-title {{ font-size: 14px; font-weight: bold; margin-bottom: 5px; color: #333; }}
            
            /* Options (Tighten gaps) */
            ul.options-list {{ margin: 0; padding-left: 20px; list-style-type: none; }}
            li {{ margin-bottom: 2px; padding: 2px 0; }}
            
            /* Answer Key */
            h4 {{ margin-top: 20px; border-bottom: 1px solid #ccc; }}
            table.answer-key {{ width: 100%; border-collapse: collapse; margin-top: 10px; font-size: 11px; }}
            th, td {{ border: 1px solid #999; padding: 5px; text-align: left; }}
            th {{ background-color: #f0f0f0; }}
        </style>
    </head>
    <body>
        <h1>Quiz: {topic}</h1>
        {clean_html}
    </body>
    </html>
    """
    
    fname = f"quiz_{uuid.uuid4().hex[:8]}.pdf"
    path = os.path.join(STATIC_FOLDER, fname)
    
    try:
        with open(path, "w+b") as f:
            pisa.CreatePDF(BytesIO(pdf_html.encode('utf-8')), dest=f)
        return jsonify({"success": True, "quiz": clean_html, "file_url": f"/static/{fname}"})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500

# --- PPT GENERATOR (FIXED: Layout) ---
@app.route('/make-ppt', methods=['POST'])
def make_ppt():
    increment_stat('text_gen')
    topic = request.form.get('topic', 'Presentation')
    src_text = request.form.get('source_text', '')
    template_file = request.files.get('template_file')
    
    prs = None
    temp_path = None

    # Load Template or Default
    try:
        if template_file and template_file.filename != '':
            temp_path = os.path.join(STATIC_FOLDER, f"temp_{uuid.uuid4().hex}.pptx")
            template_file.save(temp_path)
            prs = Presentation(temp_path)
        else:
            prs = Presentation()
    except:
        prs = Presentation()

    prompt = (
        f"Create a presentation outline for '{topic}'. Context: {src_text}.\n"
        "Output strictly 4 slides in this format:\n"
        "SLIDE: [Title of Slide]\n"
        "POINT: [Bullet point 1]\n"
        "POINT: [Bullet point 2]\n"
        "SLIDE: [Next Title]\n"
        "..."
    )
    
    ai_text = get_groq_response("You are a presentation generator.", prompt)
    if not ai_text:
        if temp_path: os.remove(temp_path)
        return jsonify({"success": False, "error": "AI Failed"})

    cleaned_text = clean_ai_text(ai_text)
    slide = None
    
    for line in cleaned_text.split('\n'):
        line = line.strip()
        if line.startswith("SLIDE:"):
            # Use layout 1 (Title + Content)
            layout_idx = 1 if len(prs.slide_layouts) > 1 else 0
            slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
            try: slide.shapes.title.text = line.replace("SLIDE:", "").strip()
            except: pass
        elif line.startswith("POINT:") and slide:
            try:
                # Add bullet to body placeholder (usually idx 1)
                text_shape = slide.placeholders[1]
                p = text_shape.text_frame.add_paragraph()
                p.text = line.replace("POINT:", "").strip()
                p.level = 0
            except: pass

    filename = f"presentation_{uuid.uuid4().hex[:8]}.pptx"
    prs.save(os.path.join(STATIC_FOLDER, filename))
    
    if temp_path and os.path.exists(temp_path):
        os.remove(temp_path)
        
    return jsonify({"success": True, "file_url": f"/static/{filename}"})

# ==============================================================================
#                               FILE TOOLS
# ==============================================================================

@app.route('/text-to-audio', methods=['POST'])
def text_to_audio():
    increment_stat('audio_gen')
    text = request.form.get('text', '')
    lang = request.form.get('target_language', 'en')
    
    if not text: return jsonify({"success": False, "error": "No text"}), 400
    
    try:
        fname = f"speech_{uuid.uuid4().hex[:8]}.mp3"
        path = os.path.join(STATIC_FOLDER, fname)
        tts = gTTS(text=text, lang=lang, slow=False)
        tts.save(path)
        return jsonify({"success": True, "file_url": f"/static/{fname}"})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500

@app.route('/audio-to-text', methods=['POST'])
def audio_to_text():
    increment_stat('transcribe')
    if 'file' not in request.files: return jsonify({"success": False, "error": "No file"}), 400
    
    file = request.files['file']
    lang = request.form.get('language', 'en-US')
    
    fname = f"temp_{uuid.uuid4().hex[:8]}.wav"
    path = os.path.join(STATIC_FOLDER, fname)
    file.save(path)
    
    try:
        r = sr.Recognizer()
        with sr.AudioFile(path) as source:
            audio_data = r.record(source)
            text = r.recognize_google(audio_data, language=lang)
        os.remove(path)
        return jsonify({"success": True, "text": text})
    except Exception as e:
        if os.path.exists(path): os.remove(path)
        return jsonify({"success": False, "error": str(e)}), 500

@app.route('/text-to-pdf', methods=['POST'])
def text_to_pdf():
    increment_stat('pdf_gen')
    html_content = request.form.get('html_content', '')
    styled_html = f"<html><body><style>body{{font-family:Helvetica;}}</style>{html_content}</body></html>"
    fname = f"doc_{uuid.uuid4().hex[:8]}.pdf"
    path = os.path.join(STATIC_FOLDER, fname)
    try:
        with open(path, "w+b") as f:
            pisa.CreatePDF(BytesIO(styled_html.encode('utf-8')), dest=f)
        return jsonify({"success": True, "file_url": f"/static/{fname}"})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500

@app.route('/video-to-audio', methods=['POST'])
def video_to_audio():
    increment_stat('vid_audio')
    if 'file' not in request.files: return jsonify({"success": False, "error": "No file"}), 400
    
    file = request.files['file']
    vid_name = f"temp_vid_{uuid.uuid4().hex[:8]}.mp4"
    vid_path = os.path.join(STATIC_FOLDER, vid_name)
    file.save(vid_path)
    
    audio_name = f"extracted_{uuid.uuid4().hex[:8]}.mp3"
    audio_path = os.path.join(STATIC_FOLDER, audio_name)
    
    try:
        with VideoFileClip(vid_path) as clip:
            clip.audio.write_audiofile(audio_path, logger=None)
        return jsonify({"success": True, "file_url": f"/static/{audio_name}"})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500
    finally:
        if os.path.exists(vid_path):
            try: os.remove(vid_path)
            except: pass

@app.route('/convert-file', methods=['POST'])
def convert_file():
    increment_stat('file_conv')
    if 'file' not in request.files: return jsonify({"success": False, "error": "No file"}), 400
    file = request.files['file']
    fmt = request.form.get('format', 'PNG').upper()
    try:
        img = PIL.Image.open(file)
        if fmt in ['JPG', 'JPEG']: img = img.convert('RGB'); fmt = 'JPEG'
        fname = f"conv_{uuid.uuid4().hex[:8]}.{fmt.lower()}"
        path = os.path.join(STATIC_FOLDER, fname)
        img.save(path, format=fmt)
        return jsonify({"success": True, "file_url": f"/static/{fname}"})
    except Exception as e: return jsonify({"success": False, "error": str(e)}), 500

@app.route('/compress-image', methods=['POST'])
def compress_image():
    increment_stat('compression')
    if 'file' not in request.files: return jsonify({"success": False, "error": "No file"}), 400
    file = request.files['file']
    try:
        img = PIL.Image.open(file).convert('RGB')
        fname = f"comp_{uuid.uuid4().hex[:8]}.jpg"
        path = os.path.join(STATIC_FOLDER, fname)
        img.save(path, "JPEG", optimize=True, quality=30)
        return jsonify({"success": True, "file_url": f"/static/{fname}"})
    except Exception as e: return jsonify({"success": False, "error": str(e)}), 500

if __name__ == '__main__':
    app.run(host='0.0.0.0', port=5000, debug=True)